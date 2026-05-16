"""Generate a single-slide PPT illustrating the 5 certification output scenarios.

Per agreed semantics:
  S1: >=30 successful runs        -> Certificate WITH statistical testing
  S2: 1..29 successful runs       -> Certificate WITHOUT statistical testing
  S3: 0 successful runs           -> NO certificate (upstream failure)
  S4: Runs OK but metrics null    -> Certificate with warnings (root cause TBD)
  S5: No fault detected (single-  -> NO certificate (single-fault fallback bucket
      fault fallback)                has no fault_category, aggregation aborts)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---- Colors ----
NAVY = RGBColor(0x1A, 0x27, 0x44)
GREY = RGBColor(0x4A, 0x55, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD7, 0xE0)

GREEN = RGBColor(0x1F, 0x7A, 0x3D);  GREEN_BG = RGBColor(0xE8, 0xF5, 0xEC)
AMBER = RGBColor(0xB0, 0x6E, 0x00);  AMBER_BG = RGBColor(0xFD, 0xF4, 0xE0)
RED   = RGBColor(0xB0, 0x2A, 0x37);  RED_BG   = RGBColor(0xFA, 0xE8, 0xEA)
ORANGE = RGBColor(0xC4, 0x4A, 0x1F); ORANGE_BG = RGBColor(0xFD, 0xEA, 0xDF)
BLUE  = RGBColor(0x0B, 0x4F, 0x9E);  BLUE_BG  = RGBColor(0xE2, 0xEE, 0xFB)


def add_text(slide, left, top, width, height, text, *, size=12, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Segoe UI"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04);  tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    box.fill.background(); box.line.fill.background()
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return box


def add_multiline(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05);  tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    box.fill.background(); box.line.fill.background()
    for i, (txt, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opts.get("space_after", 2))
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opts.get("size", 9))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", NAVY)
        r.font.name = opts.get("font", "Segoe UI")
    return box


def add_badge(slide, left, top, width, height, text, fill, color, *, size=8.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.4
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = color; shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return shp


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ---- Title ----
add_text(slide, Inches(0.3), Inches(0.18), Inches(12.7), Inches(0.45),
         "AgentCert — Certification Output Scenarios",
         size=22, bold=True, color=NAVY)
add_text(slide, Inches(0.3), Inches(0.62), Inches(12.7), Inches(0.3),
         "How the pipeline behaves under different run conditions and data availability  •  Threshold: n \u2265 30 successful runs/category for statistical testing",
         size=10, color=GREY)

# ---- Five scenario cards ----
scenarios = [
    {
        "header": GREEN, "header_bg": GREEN_BG,
        "badge": "S1", "badge_label": "CERTIFICATE + STATS",
        "title": "Full Statistical Certification",
        "condition": "\u2265 30 successful runs / category",
        "runs": "30 / 30 successful",
        "pipeline": [("Aggregation", GREEN), ("Hypothesis H-01\u2013H-09", GREEN), ("Certification", GREEN)],
        "bullets": [
            "Wilson & BCa 95% confidence intervals",
            "Kruskal\u2013Wallis, Fisher, Levene, TOST",
            "CVaR\u2089\u2085, CUSUM/EWMA drift checks",
            "All 9 hypotheses evaluated",
        ],
        "out_label": "OUTPUT",
        "out_text": "Certificate with statistical claims at 95% confidence",
        "out_fill": GREEN_BG, "out_color": GREEN,
    },
    {
        "header": AMBER, "header_bg": AMBER_BG,
        "badge": "S2", "badge_label": "CERT W/ CAVEATS",
        "title": "Aggregation-Only Report",
        "condition": "1 \u2013 29 successful runs",
        "runs": "e.g. 5 / 30 or 25 / 30",
        "pipeline": [("Aggregation", AMBER), ("Hypothesis \u2014 skipped", RED), ("Certification", AMBER)],
        "bullets": [
            "Point estimates: mean, median, P95",
            "LLM Council qualitative synthesis",
            "Directional findings only",
            "Callout: n < 30 \u2014 stats not evaluated",
        ],
        "out_label": "OUTPUT",
        "out_text": "Descriptive certificate, no statistical validation",
        "out_fill": AMBER_BG, "out_color": AMBER,
    },
    {
        "header": RED, "header_bg": RED_BG,
        "badge": "S3", "badge_label": "NO CERTIFICATE",
        "title": "No Successful Runs",
        "condition": "0 successful runs reach certifier",
        "runs": "0 / 30 successful",
        "pipeline": [("Aggregation \u2014 no input", RED), ("Hypothesis \u2014 skipped", RED), ("Certification \u2014 skipped", RED)],
        "bullets": [
            "All runs failed upstream (agent / fault injection)",
            "No metric documents produced",
            "Pipeline never triggered",
            "Action: investigate AI-engineer side",
        ],
        "out_label": "OUTPUT",
        "out_text": "No certificate \u2014 upstream failure",
        "out_fill": RED_BG, "out_color": RED,
    },
    {
        "header": ORANGE, "header_bg": ORANGE_BG,
        "badge": "S4", "badge_label": "CERT W/ WARNINGS",
        "title": "Runs OK, Metrics Null",
        "condition": "Runs successful but metrics not extracted",
        "runs": "\u2265 1 run, fields all null",
        "pipeline": [("Aggregation", ORANGE), ("Hypothesis \u2014 partial / skipped", ORANGE), ("Certification", ORANGE)],
        "bullets": [
            "Quant + qualitative metrics empty",
            "Could be trace-quality or extractor issue",
            "Certificate flags missing data prominently",
            "Root cause needs manual investigation",
        ],
        "out_label": "OUTPUT",
        "out_text": "Certificate with explicit warnings \u2014 investigate further",
        "out_fill": ORANGE_BG, "out_color": ORANGE,
    },
    {
        "header": BLUE, "header_bg": BLUE_BG,
        "badge": "S5", "badge_label": "NO CERTIFICATE",
        "title": "No Fault Detected (Single-Fault Fallback)",
        "condition": "Bucketing finds no \u201cfault: *\u201d spans",
        "runs": "Single fallback bucket created",
        "pipeline": [("Single-bucket fallback", BLUE), ("Aggregation \u2014 aborts", RED), ("Certification \u2014 skipped", RED)],
        "bullets": [
            "All events placed in one \u201cunknown\u201d bucket",
            "No fault_category, no injection timestamp",
            "Aggregator: \u201cNo fault categories found\u201d",
            "Pipeline halts before scorecard",
        ],
        "out_label": "OUTPUT",
        "out_text": "No certificate \u2014 fallback bucket lacks category",
        "out_fill": BLUE_BG, "out_color": BLUE,
    },
]

# Layout: 5 cards in a row
card_top = Inches(1.05)
card_height = Inches(5.85)
left_margin = Inches(0.3)
right_margin = Inches(0.3)
gap = Inches(0.1)
total_w = prs.slide_width - left_margin - right_margin
card_width = (total_w - gap * (len(scenarios) - 1)) / len(scenarios)

for i, sc in enumerate(scenarios):
    left = left_margin + i * (card_width + gap)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
    card.adjustments[0] = 0.03
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER; card.line.width = Pt(1)
    card.shadow.inherit = False

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_width, Inches(0.55))
    header.fill.solid(); header.fill.fore_color.rgb = sc["header_bg"]
    header.line.fill.background()

    add_badge(slide, left + Inches(0.12), card_top + Inches(0.1),
              Inches(0.42), Inches(0.32), sc["badge"], WHITE, sc["header"], size=11)
    add_text(slide, left + Inches(0.55), card_top + Inches(0.13),
             card_width - Inches(0.65), Inches(0.3),
             sc["badge_label"], size=8.5, bold=True, color=sc["header"], align=PP_ALIGN.RIGHT)

    cur = card_top + Inches(0.6)

    add_text(slide, left + Inches(0.15), cur, card_width - Inches(0.3), Inches(0.55),
             sc["title"], size=12, bold=True, color=sc["header"])
    cur += Inches(0.55)

    add_text(slide, left + Inches(0.15), cur, card_width - Inches(0.3), Inches(0.3),
             sc["condition"], size=9, bold=True, color=NAVY)
    cur += Inches(0.3)

    add_text(slide, left + Inches(0.15), cur, card_width - Inches(0.3), Inches(0.28),
             "\u25b8 " + sc["runs"], size=9, bold=True, color=sc["header"])
    cur += Inches(0.32)

    add_text(slide, left + Inches(0.15), cur, card_width - Inches(0.3), Inches(0.22),
             "Pipeline Stages", size=8, bold=True, color=GREY)
    cur += Inches(0.24)
    for label, c in sc["pipeline"]:
        add_badge(slide, left + Inches(0.15), cur,
                  card_width - Inches(0.3), Inches(0.26), label, WHITE, c, size=8)
        cur += Inches(0.3)

    cur += Inches(0.05)

    add_text(slide, left + Inches(0.15), cur, card_width - Inches(0.3), Inches(0.22),
             "What Happens", size=8, bold=True, color=GREY)
    cur += Inches(0.24)
    bullet_lines = [("\u2022  " + b, {"size": 8.5, "color": NAVY, "space_after": 3})
                    for b in sc["bullets"]]
    add_multiline(slide, left + Inches(0.15), cur,
                  card_width - Inches(0.3), Inches(1.4), bullet_lines)

    out_top = card_top + card_height - Inches(0.85)
    out_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left + Inches(0.15), out_top,
                                     card_width - Inches(0.3), Inches(0.7))
    out_box.adjustments[0] = 0.12
    out_box.fill.solid(); out_box.fill.fore_color.rgb = sc["out_fill"]
    out_box.line.color.rgb = sc["out_color"]; out_box.line.width = Pt(1)
    tf = out_box.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05);  tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = sc["out_label"]
    r1.font.size = Pt(7.5); r1.font.bold = True
    r1.font.color.rgb = sc["out_color"]; r1.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = sc["out_text"]
    r2.font.size = Pt(9); r2.font.bold = True
    r2.font.color.rgb = NAVY; r2.font.name = "Segoe UI"

add_text(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.3),
         "Statistical hypothesis testing (H-01 \u2013 H-09) is performed only in S1.  "
         "S3 and S5 do not produce a certificate.",
         size=9, bold=True, color=GREY, align=PP_ALIGN.CENTER)

out_path = r"c:\Users\shiwkumari\Projects\AgentCert\certifier\certification_scenarios.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
